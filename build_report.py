import pptx
import inflect
import dateparser
import sqlite_connector
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData, ChartData
from datetime import datetime, timedelta, date
from dateutil.relativedelta import relativedelta
from wordcloud import WordCloud

class ac_report(object):

    def __init__(self, dbfile="ac_exec_report.db"):
        self.dbfile = dbfile
        self.db = sqlite_connector.sqlite_db(dbfile)
        self.prs = pptx.Presentation(open("quarterly_report_template.pptx", "rb"))
        self.main_slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self.second_slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self.filename = f"AC Executive Report.pptx"
        # Get the date 3 months ago (the start of the quarter)
        today = datetime.combine(date.today(), datetime.min.time())
        self.start = datetime.strftime(today - relativedelta(months=3), '%Y-%m-%dT00:00:00')

    def color_cell(self, table, cell_xy, color):
        cell = table.cell(cell_xy[0], cell_xy[1])
        if color == "yellow":
            color = pptx.dml.color.RGBColor(0xFF, 0xFF, 0x00)
        elif color == "red":
            color = pptx.dml.color.RGBColor(0xFF, 0x00, 0x00)
        elif color == "orange":
            color = pptx.dml.color.RGBColor(0xFF, 0xA5, 0x00)
        elif color == "green":
            color = pptx.dml.color.RGBColor(0x00, 0xFF, 0x00)
        cell.fill.solid()
        cell.fill.fore_color.rgb = color

    def create_column_chart(self, slide, data, cs, legend=True, base_10=False):
        # Reference: data = {"category_name": {"series_name": metric}} - 
        # Reference: data = {"bengoff@gmail.com": {"logins": 45}}
        # Reference: cs = [Chart title, x, y, sizex, sizey]
        title = cs[0]
        x, y = pptx.util.Inches(cs[1]), pptx.util.Inches(cs[2])
        cx, cy = pptx.util.Inches(cs[3]), pptx.util.Inches(cs[4])
        ctype = pptx.enum.chart.XL_CHART_TYPE.COLUMN_CLUSTERED
        chart_data = CategoryChartData()
        cats = [i for i in data.keys()]
        series_names = []
        for cat in data:
            for s in data[cat]:
                if s not in series_names:
                    series_names.append(s)
        chart_data.categories = ["".join(cat.split(" -")[0]) for cat in cats]
        for series in series_names:
            chart_data.add_series("".join(series.split(" -")[0]), [data[cat].get(series) for cat in cats])

        chart = slide.shapes.add_chart(ctype, x, y, cx, cy, chart_data).chart
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = pptx.util.Pt(20)
        chart.has_legend = legend
        if legend:
            chart.legend.position = pptx.enum.chart.XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = pptx.util.Pt(10)
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = pptx.util.Pt(10)
        category_axis.has_major_gridlines = False
        red = pptx.dml.color.RGBColor(165, 36, 40)
        orange = pptx.dml.color.RGBColor(210, 93, 18)
        yellow = pptx.dml.color.RGBColor(246, 178, 33)
        greenish = pptx.dml.color.RGBColor(163, 213, 93)
        green = pptx.dml.color.RGBColor(0, 152, 68)
        # Selective coloring changes
        for x, i in enumerate(data):
            point = chart.plots[0].series[0].points[x]
            fill = point.format.fill
            fill.solid()
            if "-RED" in i:
                fill.fore_color.rgb = red
            if "-YELLOW" in i:
                fill.fore_color.rgb = yellow
            if "-GREEN" in i:
                fill.fore_color.rgb = green
        value_axis = chart.value_axis
        if base_10:
            scaling_element = value_axis._element.xpath(r"c:scaling")
            from docx.oxml import OxmlElement
            log_base = OxmlElement(r"c:logBase")
            scaling_element[0].append(log_base)
            log_base.set("val", "10")
        value_axis.tick_labels.font.size = pptx.util.Pt(10)

    def create_pie_chart(self, data, cs):
        title = cs[0]
        x, y = Inches(cs[1]), Inches(cs[2])
        cx, cy = Inches(cs[3]), Inches(cs[4])
        ctype = pptx.enum.chart.XL_CHART_TYPE.PIE
        chart_data = CategoryChartData()
        chart_data.categories = [i[0] for i in data]
        chart_data.add_series("Count", [i[1] for i in data])
        chart = self.main_slide.shapes.add_chart(ctype, x, y, cx, cy, chart_data).chart
        chart.has_legend = False
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = pptx.util.Pt(20)
        chart.has_title = False
        if chart.has_legend:
            chart.legend.include_in_layout = False
            chart.legend.position = pptx.enum.chart.XL_LEGEND_POSITION.BOTTOM
            chart.legend.font.size = pptx.util.Pt(12)
        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.font.size = pptx.util.Pt(12)
        data_labels = chart.plots[0].data_labels
        data_labels.position = pptx.enum.chart.XL_LABEL_POSITION.BEST_FIT
        data_labels.number_format = '0%'
        for series in chart.series:
            for x, point in enumerate(series.points):
                point.data_label.text_frame.text = data[x][0]
                point.data_label.font.size = pptx.util.Pt(12)
                for para in point.data_label.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(8)
                        run.font.bold = True

    def create_table(self, slide, data, ts):
        ''' create table with slide, [data], and [table settings] '''

        def setTblFont(table, pts, wids):
            if wids == True:
                table.columns[0].width = pptx.util.Inches(1.5)
                table.columns[1].width = pptx.util.Inches(0.5)
                table.columns[2].width = pptx.util.Inches(0.5)
                table.columns[3].width = pptx.util.Inches(0.5)
                table.columns[4].width = pptx.util.Inches(1)

            def iter_cells(table):
                for row in table.rows:
                    for cell in row.cells:
                        yield cell

            for cell in iter_cells(table):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = pptx.util.Pt(pts)
                        run.font.color.rgb = pptx.dml.color.RGBColor(0x3f, 0x2c, 0x36)
            for x, row in enumerate(table.rows):
                if x > 0: break
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = pptx.dml.color.RGBColor(211, 211, 211)

        rows = len(data)
        cols = len(data[0])
        x, y = Inches(ts[0]), Inches(ts[1])
        cx, cy = Inches(ts[2]), Inches(ts[3])
        fontSz = ts[4]
        colWids = ts[5]

        table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
        for x, row in enumerate(data):
            for y in range(len(data[x])):
                if row[y] and not isinstance(row[y], list) and (isinstance(row[y], int) or isinstance(row[y], float) or row[y].isnumeric()):
                    table.cell(x,y).text = f"{int(row[y]):,}"
                else:
                    cell = table.cell(x, y)
                    #table.cell(x,y).text = str(row[y])
                    if isinstance(row[y], list):
                        text_str, url = row[y][0], row[y][1]
                        cell_link = cell.text_frame.paragraphs[0].add_run()
                        cell_link.text = text_str
                        hlink = cell_link.hyperlink
                        hlink.address = url
                    else:
                        cell.text = str(row[y])
        for x, wid in enumerate(colWids):
            table.columns[x].width = pptx.util.Inches(wid)
        setTblFont(table, fontSz, False)
        return table

    def create_line_chart(self, data, cs, legend=True):
        """ Data here would look like when not dict:
        [
        [["date", series_name_1], ["2020-10-01", 100], ["2020-10-02", 90]],
        [["date", series_name_2], ["2020-10-01", 45], ["2020-10-02", 213]]
        ]
        When a dict:
        {
        series1_name: [[category1, value], [category2, value]],
        series2_name: [[category1, value], [category2, value]]
        }
        """
        title = cs[0]
        x, y = Inches(cs[1]), Inches(cs[2])
        cx, cy = Inches(cs[3]), Inches(cs[4])
        ctype = pptx.enum.chart.XL_CHART_TYPE.LINE
        chart_data = ChartData()
        if isinstance(data, dict):
            for series in data:
                chart_data.categories = [i[0] for i in data[series]]
                chart_data.add_series(series.title(), [i[1] for i in data[series]])
        elif isinstance(data, list):
            for series in data:
                chart_data.categories = [i[0] for i in series[1:]]
                chart_data.add_series(series[0][1].title(), [i[1] for i in series[1:]])
        chart = self.main_slide.shapes.add_chart(ctype, x, y, cx, cy, chart_data).chart
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = pptx.util.Pt(20)
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = pptx.util.Pt(10)
        category_axis.has_major_gridlines = False
        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = pptx.util.Pt(10)
        value_axis.has_major_gridlines = False
        value_axis.minimum_scale = 0
        if legend:
            chart.legend.include_in_layout = True
            chart.legend.position = pptx.enum.chart.XL_LEGEND_POSITION.TOP
            chart.legend.font.size = pptx.util.Pt(14)
            if len(data) > 2:
                chart.legend.font.size = pptx.util.Pt(10)
            if len(data) > 4:
                chart.legend.font.size = pptx.util.Pt(8)
        else:
            chart.has_legend = False
            #chart.plots[0].has_data_labels = True
            #data_labels = chart.plots[0].data_labels
            #data_labels.position = pptx.enum.chart.XL_LABEL_POSITION.BEST_FIT
        plot = chart.plots[0]
        series = plot.series
        for i in series:
            i.smooth = False
            i.format.line.width = Pt(8)
            if len(series) > 2:
                i.format.line.width = Pt(4)
                i.smooth = True
            if i.name == "Warning":
                line = i.format.line.color.rgb = pptx.dml.color.RGBColor(255,255,0)
            if i.name == "Danger":
                line = i.format.line.color.rgb = pptx.dml.color.RGBColor(255,0,0)
        return chart

    def create_stacked_bar_chart(self, data, cs):
        chart_data = CategoryChartData()
        chart_data.categories = [i[0] for i in data[1:]]
        for x in range(1, len(data[0])):
            chart_data.add_series(data[0][x], [i[x] for i in data[1:]])
        x, y, cx, cy = Inches(.81), Inches(12.64), Inches(6.35), Inches(2.36)
        ctype = pptx.enum.chart.XL_CHART_TYPE.BAR_STACKED
        graphic_frame = self.main_slide.shapes.add_chart(ctype, x, y, cx, cy, chart_data)
        title = "Sensor Support Levels"
        chart = graphic_frame.chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = pptx.enum.chart.XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = pptx.util.Pt(14)
        chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = pptx.util.Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = pptx.dml.color.RGBColor(0xFF, 0xFF, 0xFF)
        chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = pptx.util.Pt(12)
        chart.value_axis.maximum_scale = 100
        chart.category_axis.tick_labels.font.size= pptx.util.Pt(12)
        chart.value_axis.tick_labels.font.size= pptx.util.Pt(12)

    def create_stacked_bar_chart_2(self, cats, data, cs):
        chart_data = CategoryChartData()
        chart_data.categories = cats
        for x, row in enumerate(data):
            chart_data.add_series("", row)
        x, y = Inches(cs[1]), Inches(cs[2])
        cx, cy = Inches(cs[3]), Inches(cs[4])
        ctype = pptx.enum.chart.XL_CHART_TYPE.BAR_STACKED
        graphic_frame = self.main_slide.shapes.add_chart(ctype, x, y, cx, cy, chart_data)
        chart = graphic_frame.chart
        title = cs[0]
        chart.chart_title.text_frame.text = title
        chart.has_title = True
        chart.has_legend = False
        #chart.legend.position = pptx.enum.chart.XL_LEGEND_POSITION.BOTTOM
        #chart.legend.include_in_layout = False
        #chart.legend.font.size = pptx.util.Pt(14)
        chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = pptx.util.Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = pptx.dml.color.RGBColor(0xFF, 0xFF, 0xFF)
        chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = pptx.util.Pt(12)
        #chart.value_axis.manimum_scale = 0
        #chart.value_axis.maximum_scale = cs[5]
        chart.category_axis.tick_labels.font.size= pptx.util.Pt(12)
        chart.value_axis.tick_labels.font.size= pptx.util.Pt(12)

    def create_stacked_column_chart(self, cats, data, cs, legend_pos):
        leg_pos = {
            "bottom": pptx.enum.chart.XL_LEGEND_POSITION.BOTTOM,
            "left": pptx.enum.chart.XL_LEGEND_POSITION.LEFT,
            "right":pptx.enum.chart.XL_LEGEND_POSITION.RIGHT
        }
        chart_data = CategoryChartData()
        chart_data.categories = cats
        for x, row in enumerate(data):
            chart_data.add_series(row[0], row[1:])
        x, y = Inches(cs[1]), Inches(cs[2])
        cx, cy = Inches(cs[3]), Inches(cs[4])
        ctype = pptx.enum.chart.XL_CHART_TYPE.COLUMN_STACKED
        graphic_frame = self.main_slide.shapes.add_chart(ctype, x, y, cx, cy, chart_data)
        chart = graphic_frame.chart
        title = cs[0]
        chart.chart_title.text_frame.text = title
        chart.has_title = True
        chart.has_legend = True
        chart.legend.position = leg_pos[legend_pos]
        chart.legend.include_in_layout = False
        chart.legend.font.size = pptx.util.Pt(12)
        if len(data) == 1:
            chart.has_legend = False
        chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = pptx.util.Pt(10)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = pptx.dml.color.RGBColor(0xFF, 0xFF, 0xFF)
        chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = pptx.util.Pt(12)
        #chart.value_axis.manimum_scale = 0
        #chart.value_axis.maximum_scale = cs[5]
        chart.category_axis.tick_labels.font.size= pptx.util.Pt(12)
        chart.value_axis.tick_labels.font.size= pptx.util.Pt(12)

    def create_text_box(self, slide, text, ts, color=False, alignment=False):
        x, y = Inches(ts[0]), Inches(ts[1])
        tx, ty = Inches(ts[2]), Inches(ts[3])
        text_placeholder = slide.shapes.add_textbox(x, y, tx, ty)
        text_frame = text_placeholder.text_frame
        text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{text} "
        run.font.size = Pt(ts[4])
        run.font.bold = True
        if color:
            run.font.color.rgb = color
        if alignment == "center":
            p.alignment = pptx.enum.text.PP_ALIGN.CENTER
        elif alignment == "right":
            p.alignment = pptx.enum.text.PP_ALIGN.RIGHT
        return text_frame

    def create_image(self, image_path, settings):
        x, y = Inches(settings[0]), Inches(settings[1])
        ix, iy = Inches(settings[2]), Inches(settings[3])
        pic = self.main_slide.shapes.add_picture(image_path, x, y, ix, iy)
        return pic

    def create_rectangle(self, slide, ss):
        # ss = (x, y, sx, sy, line_width)
        x, y = Inches(ss[0]), Inches(ss[1])
        sx, sy = Inches(ss[2]), Inches(ss[3])
        ld = Pt(ss[4])
        rr = pptx.enum.shapes.MSO_SHAPE.ROUNDED_RECTANGLE
        shapes = slide.shapes
        shape = shapes.add_shape(rr, x, y, sx, sy)
        fill = shape.fill
        fill.background()
        shape.line.width = ld

    def create_arrow(self, slide, ss, text=None):
        # ss = (x, y, sx, sy, line_width, direction, rotation)
        x, y = Inches(ss[0]), Inches(ss[1])
        sx, sy = Inches(ss[2]), Inches(ss[3])
        ld = Pt(ss[4])
        dir_lookup = {
            "up": pptx.enum.shapes.MSO_SHAPE.UP_ARROW, 
            "down": pptx.enum.shapes.MSO_SHAPE.DOWN_ARROW,
            "left": pptx.enum.shapes.MSO_SHAPE.LEFT_ARROW,
            "right": pptx.enum.shapes.MSO_SHAPE.RIGHT_ARROW,
        }
        arrow = dir_lookup[ss[5]]
        shape = slide.shapes.add_shape(arrow, x, y, sx, sy)
        if text:
            text_frame = shape.text_frame
            text_frame._bodyPr.attrib.update({'vert': 'vert'})
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.font.color.rgb = pptx.dml.color.RGBColor(242, 242, 242)
            run.text = f"{text} "
        shape.rotation = ss[6]

    def add_divider(self, slide, loc, div_text=None):
        shapes = slide.shapes
        left, width, height = Inches(0), Inches(13.31), Inches(.42)
        top = Inches(loc)
        shape = shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, left, top, width, height)
        text_placeholder = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_placeholder.text_frame
        p = text_frame.paragraphs[0]
        p.alignment = pptx.enum.text.PP_ALIGN.CENTER
        run = p.add_run()
        run.text = div_text
        run.font.size = Pt(18)

    def add_cross(self, slide, loc):
        shapes = slide.shapes
        left, width, height = Inches(0), Inches(5), Inches(.42)
        top = Inches(loc)
        vert = shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, left, top, width, height)

        left, width, height = Inches(4), Inches(.42), Inches(5)
        top = Inches(loc)
        horiz = shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, left, top, width, height)  

    def delete_index(self, slide, idx):
        textbox = slide.placeholders[idx]
        sp = textbox.element
        sp.getparent().remove(sp)

    def delete_slide(self, idx):
        slide_id = self.prs.slides._sldIdLst[idx].rId
        self.prs.part.drop_rel(slide_id)
        del self.prs.slides._sldIdLst[idx]

    def add_title(self, slide, text):
        title = slide.shapes.title
        p = title.text_frame.paragraphs[0]
        p.alignment = pptx.enum.text.PP_ALIGN.CENTER
        title.text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
        run = p.add_run()
        run.text = text
        run.font.size = Pt(28)
        font = run.font
        run = p.add_run()
        today = datetime.strftime(date.today(), "%b %d, %Y")
        run.text = f"{today}"
        run.font.size = Pt(20)

    def add_deployment_chart(self):
        query = f"""
        select
        case enforcementLevel
        when "20" then "High"
        when "30" then "Medium"
        when "35" then "Local Approval"
        when "40" then "Low"
        when "60" then "None (Visibility)"
        when "80" then "None (Disabled)"
        else enforcementLevel
        end as deployment_level,
        sum(case when lastPollDate >= (select max(date(lastPollDate, '-30 days')) from computer) THEN 1 else 0 end),
        sum(case when lastPollDate >= (select max(date(lastPollDate, '-7 days')) from computer) THEN 1 else 0 end),
        sum(case when lastPollDate >= (select max(date(lastPollDate, '-1 days')) from computer) THEN 1 else 0 end)
        from computer
        where enforcementLevel <> 0
        and lastPollDate >= (select max(date(lastPollDate, '-30 days')) from computer)
        group by deployment_level
        order by enforcementLevel;
        """
        data = self.db.query_data(query)
        cats = ("30 Day Checkin", "7 Day Checkin", "1 Day Checkin")
        cs = ["Deployment", .58, 1.82, 4, 4]
        self.create_stacked_column_chart(cats, data, cs, "bottom")

    def add_deployment_text(self):
        query = """
        select count(*)
        from computer
        where enforcementLevel <> 0
        and lastPollDate >= (select max(date(lastPollDate, '-30 days')) from computer);
        """
        total_deployed = "{:,}".format(self.db.query_data(query)[0][0])
        total_deployed = f"{total_deployed} Total Deployed"
        ts = [.32, 5.65, 1, 1, 32]
        self.create_text_box(self.main_slide, total_deployed, ts)

    def add_versions_chart(self):
        query = f"""
        select c.agentVersion || '(' || substr(c.osShortName, 0, 4) || ')' || " -" || dvs.support_level as version,
        count(*)
        from computer c
        left join deployed_version_status dvs on
            c.agentVersion = dvs.version and
            lower(substr(c.osShortName, 0, 4)) = lower(dvs.os)
        where c.lastPollDate >= (select max(date(lastPollDate, '-30 days')) from computer)
        and c.enforcementLevel <> 0
        group by c.agentVersion || '(' || substr(c.osShortName, 0, 4) || ')' || " -" || dvs.support_level;
        """
        data = self.db.query_data(query)
        reds, yellows, greens = [], [], []
        for row in data:
            row[0] = row[0].replace("-EOL", "-RED").replace("-EX", "-YELLOW").replace("-ST", "-GREEN")
            if "RED" in row[0]: reds.append(row)
            elif "YELLOW" in row[0]: yellows.append(row)
            elif "GREEN" in row[0]: greens.append(row)
        reds = sorted(reds, key=lambda x: x[0])
        yellows = sorted(yellows, key=lambda x: x[0])
        greens = sorted(greens, key=lambda x: x[0])
        sorted_data = reds + yellows + greens
        sd_dict = {}
        for x, v in enumerate([i[0] for i in sorted_data]):
            sd_dict[v] = {"count": sorted_data[x][1]}
        cs = ["Versions", 4.9, 1.7, 3.68, 3.68]
        if sorted_data:
            self.create_column_chart(self.main_slide, sd_dict, cs, legend=False, base_10=True)

    def add_support_level_table(self):
        query = f"""
        select 
        case dvs.support_level
        when 'ST' then 'Standard (Green)'
        when 'EX' then 'Extended (Yellow)'
        when 'EOL' then 'EOL (Red)'
        end as level,
        count(*)
        from computer c
        left join deployed_version_status dvs on 
            c.agentVersion = dvs.version and 
            lower(substr(c.osShortName, 0, 4)) = lower(dvs.os)
        where c.lastPollDate >= (select max(date(lastPollDate, '-30 days')) from computer)
        and c.enforcementLevel <> 0
        group by dvs.support_level;
        """
        data = self.db.query_data(query)
        total = sum([i[1] for i in data])
        for row in data:
            perc = round(row[1]/total * 100,2)
            row.append(f"{perc}%")
        header = ("Level", "Count", "Percent")
        data.insert(0, header)
        ts = [5.7, 5.51, 1.29, 1.07, 8, (1.25, .65, .65)]
        self.create_table(self.main_slide, data, ts)

    def add_la_over_time_chart(self):
        query = """
        select timestamp
        from block_events;
        """
        date_fmt = "%Y-%m-%dT%H:%M:%Sz"
        data = [datetime.strptime(i[0], date_fmt) for i in self.db.query_data(query)]
        data.sort()
        start, end = data[0].date(), datetime.utcnow().date()
        data = [datetime.strftime(i, "%b-%d") for i in data]
        all_days = {datetime.strftime(k, "%b-%d"):0 for k in [start + timedelta(days=x) for x in range((end-start).days)]}
        for i in data:
            all_days[i]+=1
        all_days = [[[k, v] for k,v in all_days.items()]]
        all_days[0].insert(0, ["Date", "Count"])
        cs = ["Local Approval Usage", 8.71, 1.71, 4.5, 4.11]
        self.create_line_chart(all_days, cs, legend=False)

    def add_undesired_enforcement_text(self):
        query = f"""
        select
        case when enforcementLevel in ("35", "40") then "Low / Local Approval"
        when enforcementLevel in ("60", "80") then "Visibility / Disabled"
        end as deployment_level,
        sum(case when lastPollDate >= (select max(date(lastPollDate, '-30 days')) from computer) THEN 1 else 0 end) as count
        from computer
        where enforcementLevel in ("35", "40", "60", "80")
        group by deployment_level
        order by enforcementLevel;
        """
        data = self.db.query_data(query)
        red = pptx.dml.color.RGBColor(165, 36, 40)
        bt = ""
        counts = [str(row[1]) for row in data]
        cats = [f" in {row[0]}" for row in data]

        count_text_box = "\n".join(counts)
        ts = [8.2, 5.87, 1, .76, 24]
        self.create_text_box(self.main_slide, count_text_box, ts, color=red, alignment="right")

        cat_text_box = "\n".join(cats)
        ts = [9.01, 5.87, 4, .76, 24]
        self.create_text_box(self.main_slide, cat_text_box, ts)

    def add_logins_chart(self):
        query = """
        select username, 
        count(*)
        from console_logins
        where timestamp >= (select max(date(timestamp, '-90 days')) from console_logins)
        """
        data = self.db.query_data(query)
        data_dict = {i[0]: {"logins": i[1]} for i in data}
        cs = ["Logins (last 90d)", .21, 7.41, 5.04, 4.68]
        self.create_column_chart(self.main_slide, data_dict, cs, legend=False)

    def add_policy_table(self):
        query = """
		select p.name,
		case p.enforcementLevel
			when "20" then "High"
			when "30" then "Medium"
			when "35" then "Local Approval"
			when "40" then "Low"
			when "60" then "None (Visibility)"
			when "80" then "None (Disabled)"
		end as enforcement_level,
		count(c.id)
		from policy p
		left join computer c on p.id = c.policyId
		group by p.name, p.enforcementlevel
		order by count(c.id) desc;
        """
        data = self.db.query_data(query)
        for x, i in enumerate(data):
            if len(i[0]) > 27:
                data[x][0] = f"{i[0][:27]}..."
        all_others = sum([i[2] for i in data[14:]])
        data = data[:14]
        data.append(["-All Others-", "NA", all_others])
        reds, yellows = [], []
        for x, row in enumerate(data):
            if row[1] in ("Low", "None (Visibility)", "None (Disabled)"):
                for y in range(3): yellows.append([x+1, y])
            for level in ("Low", "Medium", "High"):
                if level.lower() in row[0].lower() and row[1] != level:
                    for y in range(3): reds.append([x+1, y])
        header = ("Policy", "Enforcement", "Endpoints")
        data.insert(0, header)
        ts = [8.96, 7.95, 1.29, 1.07, 9, (2.2, 1.2, .85)]
        table = self.create_table(self.main_slide, data, ts)
        for i in yellows: self.color_cell(table, i, "yellow")
        for i in reds: self.color_cell(table, i, "red")
        self.create_text_box(self.main_slide, "Policies", [10.34, 7.42, 1.5, .34, 20])

    def add_pruning_text_metrics(self):
        query = "select name, value from serverConfig where name in ('PurgeEventLogPeriod', 'PurgeEventThreshold')"
        purges = {k:v for k, v in self.db.query_data(query)}
        query = "select receivedtimestamp from oldest_event;"
        purges["oldest_event"] = self.db.query_data(query)[0][0]
        query = "select receivedtimestamp from newest_event;"
        purges["newest_event"] = self.db.query_data(query)[0][0]
        query = "select count from event_count_30d"
        purges["actual_events"] = "{:,}".format(self.db.query_data(query)[0][0])
        purges["days"] = int(int(purges["PurgeEventLogPeriod"]) / (60 * 60 * 24))
        words = inflect.engine()
        purges["events"] = words.number_to_words(purges["PurgeEventThreshold"]).title()
        dt_format = "%Y-%m-%dT%H:%M:%S.%fZ"
        #newest = datetime.strptime(purges["newest_event"].upper(), dt_format)
        #oldest = datetime.strptime(purges["oldest_event"].upper(), dt_format)
        newest = dateparser.parse(purges["newest_event"])
        oldest = dateparser.parse(purges["oldest_event"])
        purges["retention"] = (newest - oldest).days
        red = pptx.dml.color.RGBColor(0xFF, 0x00, 0x00)
        green = pptx.dml.color.RGBColor(0x00, 0xFF, 0x00)
        color = red if purges["retention"] < purges["days"] else green
        text = f"Retention Settings:\n{purges['days']} Days\n{purges['events']} Events"
        ts = [6.48, 8.55, 1, 1, 24]
        self.create_text_box(self.main_slide, text, ts, alignment="center")
        text = f"Actual Retention:\n{purges['retention']} Days\n{purges['actual_events']} Events"
        ts = [6.48, 9.85, 1, 1, 24]
        self.create_text_box(self.main_slide, text, ts, color=color, alignment="center")

    def add_blocks_over_time_chart(self):
        data_dict = {}
        date_fmt = "%Y-%m-%dT%H:%M:%Sz"
        days_back = -60

        query = f"""
        select timestamp
        from block_events
        --where timestamp >= (select max(date(timestamp, '{days_back} days')) from block_events);
        where timestamp >= date('now', '{days_back} days');
        """
        data = [dateparser.parse(i[0]) for i in self.db.query_data(query)]
        data.sort()
        data_dict["blocks"] = data

        query = f"""
        select dateCreated
        from approvalRequest
        where requestType = 1
        --and dateCreated >= (select max(date(dateCreated, '{days_back} days')) from approvalRequest);
        and dateCreated >= date('now', '{days_back} days');

        """
        data = [dateparser.parse(i[0]) for i in self.db.query_data(query)]
        data.sort()
        data_dict["requests"] = data

        query = f"""
        select dateModified
        from approvalRequest
        where resolution >1
        --and dateModified > (select max(date(dateModified, '{days_back} days')) from approvalRequest);
        and dateModified >= date('now', '{days_back} days');
        """
        data = [dateparser.parse(i[0]) for i in self.db.query_data(query)]
        data.sort()
        data_dict["approved"] = data

        query = f"""
        select dateModified
        from approvalRequest
        where resolution = 1
        --and dateModified > (select max(date(dateModified, '{days_back} days')) from approvalRequest);
        and dateModified >= date('now', '{days_back} days');
        """
        data = [dateparser.parse(i[0]) for i in self.db.query_data(query)]
        data.sort()
        data_dict["rejected"] = data

        start = min([data_dict[metric][0].date() for metric in data_dict if data_dict[metric]])
        start = datetime.date(datetime.now()) - timedelta(days=60) 

        end = datetime.utcnow().date()
        span = (end - start).days + 1

        # initialize to zero out all days and metrics
        all_days = {}
        for metric in data_dict:
            #all_days[metric] = {str(k):0 for k in [start + timedelta(days=x) for x in range(span)]}
            all_days[metric] = {k:0 for k in [start + timedelta(days=x) for x in range(span)]}

        # add up the actual counts
        for metric in data_dict:
            for i in data_dict[metric]:
                #date = str(i.date())
                date = i.date()
                all_days[metric][date] += 1

        # format the dates
        all_days = {metric: {datetime.strftime(i, "%b-%d"):all_days[metric][i] for i in all_days[metric]} for metric in all_days}

        rows = []
        import random
        for x, metric in enumerate(all_days):
            rows.append([["Date", metric]])
            for date in all_days[metric]:
                rows[x].append([date, all_days[metric][date]])
                #rows[x].append([date, random.randint(0, 20)])
        cs = ["Blocks Over Time", .4, 12.83, 12.61, 3.32]
        self.create_line_chart(rows, cs, legend=True)

    def add_top_blocking_files_table(self):
        days_back = "30"
        query = f"""
        select filename,
        count(*)
        from block_events
        --where timestamp >= (select max(date(timestamp, '-30 days')) from block_events)
        where timestamp >= date('now', '-{days_back} days')
        group by filename
        order by count(*) desc;
        """
        data = self.db.query_data(query)[:12]
        data.insert(0, ("File", "Blocks"))
        ts = [.12, 16.34, 1.29, 1.07, 9, (3, 1)]
        self.create_table(self.main_slide, data, ts)

    def add_block_text(self):
        ts = [4.58, 16.11, 4.27, 1, 18]
        days_back = "30"
        text = f"Over the last {days_back} days there were:"
        self.create_text_box(self.main_slide, text, ts, alignment="center")

        query = f"select count(*) from approvalRequest where dateModified >= date('now', '-{days_back} days');"
        total_approvals = self.db.query_data(query)[0][0]

        query = f"""
        select
        count(*)
        from block_events
        --where timestamp >= (select max(date(timestamp, '-90 days')) from block_events);
        where timestamp >= date('now', '-{days_back} days');
        """
        total_blocks = self.db.query_data(query)[0][0]

        ts = [5.72, 16.81, 2, 1, 40]
        self.create_text_box(self.main_slide, total_blocks, ts, alignment="center")
        ts[1] = 18.22
        self.create_text_box(self.main_slide, total_approvals, ts, alignment="center")

        ts = [5.72, 17.39, 2, 1, 18]
        text = "Total Blocks"
        self.create_text_box(self.main_slide, text, ts, alignment="center")
        text = "Approval Requests Submitted"
        ts[1] = 18.8
        self.create_text_box(self.main_slide, text, ts, alignment="center")

    def add_approval_breakdown_chart(self):
        days_back = "30"
        query = f"""
        select
        case ar.resolution
            when "0" then 'Not Resolved'
            when "1" then 'Rejected'
            when "2" then 'Approved'
            when "3" then 'Rule Change'
            when "4" then 'Installer'
            when "5" then 'Updater'
            when "6" then 'Publisher'
            when "7" then 'Other'
        end as resolution_type,
        count(*)
        from approvalRequest ar
        --where dateCreated >= (select max(date(timestamp, '-30 days')) from block_events)
        where dateCreated >= date('now', '-{days_back} days')
        group by ar.resolution;
        """
        data = self.db.query_data(query)
        if data:
            cs = ["Approval Request Resolution", 9.32, 15.87, 4.16, 4.16]
            self.create_pie_chart(data, cs)

    def add_arrows(self):
        arrow_s = [4.3, 17.01, 1.5, .5, 10, "left", 0]
        self.create_arrow(self.main_slide, arrow_s)

        arrow_s = [7.15, 18.47, 2.15, .5, 10, "right", 0]
        self.create_arrow(self.main_slide, arrow_s)


def create_report():
    db = sqlite_connector.sqlite_db("ac_exec_report.db")
    report = ac_report()
    report.add_title(report.main_slide, "CB App Control Executive Report\n")
    report.add_divider(report.main_slide, 1.2, "Agent Overview")
    report.add_deployment_chart()
    report.add_deployment_text()
    report.add_versions_chart()
    report.add_support_level_table()
    report.add_la_over_time_chart()
    report.add_undesired_enforcement_text()

    report.add_divider(report.main_slide, 6.82, "Console Access & Environment Configuration")
    report.add_logins_chart()
    report.add_pruning_text_metrics()
    report.add_policy_table()

    report.add_divider(report.main_slide, 12.25, "Blocks and Approvals")
    report.add_blocks_over_time_chart()
    report.add_top_blocking_files_table()
    report.add_block_text()
    report.add_arrows()
    report.add_approval_breakdown_chart()

    report.delete_slide(1)
    report.prs.save(report.filename)

if __name__ == "__main__":
    create_report()
